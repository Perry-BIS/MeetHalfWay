"""Build an editable PPTX with the two MeetHalfway demo screenshots
and RAGTrip-style annotations (dashed orange callouts + connector lines
+ bold orange labels).

Output: figures/demo_screenshots.pptx (widescreen 13.33 x 7.5 in).
"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


FIG = Path(__file__).resolve().parent
COVER_PNG  = FIG / "cover.png"
RESULT_PNG = FIG / "result_map.png"
OUT_PPTX   = FIG / "demo_screenshots.pptx"

# ---------- palette ----------
ORANGE    = RGBColor(0xE1, 0x5A, 0x1E)   # primary annotation color
DARK_GRAY = RGBColor(0x36, 0x45, 0x4F)   # slide titles
SUB_GRAY  = RGBColor(0x55, 0x55, 0x55)   # caption subtext

# ---------- image area on each slide ----------
SLIDE_W, SLIDE_H = 13.333, 7.5            # inches (widescreen)
IMG_X, IMG_Y     = 0.40, 0.95             # inches
IMG_W, IMG_H     = 7.50, 5.36             # inches  (2800x2000 native -> 1.4:1)
PX_W, PX_H       = 2800, 2000             # screenshot native pixels

LABEL_X = 8.10                            # inches; right column for labels
LABEL_W = 5.00                            # inches


def px_to_in_x(px):  return IMG_X + (px / PX_W) * IMG_W
def px_to_in_y(py):  return IMG_Y + (py / PX_H) * IMG_H


def _set_dash(shape, dash="dash"):
    """Force a preset dash on the shape's line element via raw XML."""
    ln = shape.line._get_or_add_ln()
    for old in ln.findall(qn("a:prstDash")):
        ln.remove(old)
    pd = etree.SubElement(ln, qn("a:prstDash"))
    pd.set("val", dash)


def add_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.40), Inches(0.20),
                                  Inches(12.5), Inches(0.55))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = DARK_GRAY
    r.font.name = "Calibri"


def add_callout_box(slide, px1, py1, px2, py2):
    x1, y1 = px_to_in_x(px1), px_to_in_y(py1)
    x2, y2 = px_to_in_x(px2), px_to_in_y(py2)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x1), Inches(y1),
        Inches(x2 - x1), Inches(y2 - y1),
    )
    shp.fill.background()
    shp.line.color.rgb = ORANGE
    shp.line.width = Pt(2.0)
    _set_dash(shp, "dash")
    # Remove default shape text padding so empty rectangle stays clean
    return shp


def add_label(slide, y_in, headline, subtext=None):
    h = 0.95 if subtext else 0.5
    tb = slide.shapes.add_textbox(
        Inches(LABEL_X), Inches(y_in),
        Inches(LABEL_W), Inches(h),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = headline
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = ORANGE
    r.font.name = "Calibri"

    if subtext:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtext
        r2.font.size = Pt(11)
        r2.font.italic = True
        r2.font.color.rgb = SUB_GRAY
        r2.font.name = "Calibri"
    return tb


def add_line(slide, box, label):
    """Connector from box's right-mid -> label's left-mid."""
    x1 = box.left + box.width
    y1 = box.top + box.height // 2
    x2 = label.left
    y2 = label.top + label.height // 2
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = ORANGE
    conn.line.width = Pt(1.5)
    return conn


# ============================================================
prs = Presentation()
prs.slide_width  = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]


# ---------- Slide 1: cover / landing ----------
s1 = prs.slides.add_slide(BLANK)
add_title(s1, "MeetHalfway Demo Interface  —  Home Page")
s1.shapes.add_picture(str(COVER_PNG),
                      Inches(IMG_X), Inches(IMG_Y),
                      Inches(IMG_W), Inches(IMG_H))

# (px1, py1, px2, py2, label_y_inch, headline, subtext)
ANNS_S1 = [
    (40, 30, 660, 500, 1.05,
     "Project title and positioning",
     "Project name and tagline frame the system as a "
     "privacy-first meeting-place planner."),
    (790, 30, 1280, 500, 2.20,
     "Design principles at a glance",
     "Core goal, privacy policy, recommendation target, "
     "and interaction model presented as summary cards."),
    (40, 540, 1300, 1090, 3.35,
     "Six-step privacy-separated workflow",
     "From private check-in to commute radius, overlap "
     "search, voting, time alignment, and final suggestion."),
    (40, 1100, 1300, 1450, 4.50,
     "Privacy guarantees",
     "Participants never exchange raw coordinates with "
     "one another; only outcomes are shown."),
    (1080, 1470, 1300, 1580, 5.65,
     "Entry point into the meeting query",
     "The Next button advances the user from the landing "
     "page to the participant input flow."),
]

for px1, py1, px2, py2, ly, head, sub in ANNS_S1:
    b = add_callout_box(s1, px1, py1, px2, py2)
    l = add_label(s1, ly, head, sub)
    add_line(s1, b, l)


# ---------- Slide 2: search result map ----------
s2 = prs.slides.add_slide(BLANK)
add_title(s2, "MeetHalfway Demo Result  —  Shared Feasible Region & Candidates")
s2.shapes.add_picture(str(RESULT_PNG),
                      Inches(IMG_X), Inches(IMG_Y),
                      Inches(IMG_W), Inches(IMG_H))

# Annotations ordered top-down on the right column to minimise line-crossing.
ANNS_S2 = [
    # Polygon: a slice at the top-left of the green shaded area
    (350, 150, 850, 450, 1.05,
     "Shared feasible region",
     "Green shaded polygon is the intersection of the two "
     "travel-time isochrones — the jointly reachable area."),
    # Candidate cluster (green pins in upper-middle)
    (1380, 380, 1920, 650, 2.20,
     "Ranked candidate venues",
     "Green pins denote candidate POIs retrieved inside "
     "the shared region and ranked by the fairness score."),
    # Participant B (blue pin, right side)
    (2140, 470, 2270, 620, 3.35,
     "Participant B starting location",
     "Blue pin marks the second participant's "
     "geocoded starting point."),
    # Midpoint (pink/magenta pin near centre)
    (1280, 800, 1410, 980, 4.50,
     "Geometric midpoint baseline",
     "Magenta pin shows the naive midpoint a non-"
     "fairness-aware tool would return, for contrast."),
    # Participant A (red pin, lower-left)
    (170, 1100, 300, 1270, 5.65,
     "Participant A starting location",
     "Red pin marks the first participant's "
     "geocoded starting point."),
]

for px1, py1, px2, py2, ly, head, sub in ANNS_S2:
    b = add_callout_box(s2, px1, py1, px2, py2)
    l = add_label(s2, ly, head, sub)
    add_line(s2, b, l)


prs.save(OUT_PPTX)
print(f"saved: {OUT_PPTX}")
